import re
from pptx import Presentation
import json
import os
import glob

def clean_and_tokenize(text):
    """Limpia y tokeniza texto"""
    text = text.lower()
    # Normaliza acentos para español
    text = text.replace('á','a').replace('é','e').replace('í','i').replace('ó','o').replace('ú','u')
    text = re.sub(r'[^\w\s]', ' ', text)
    text = re.sub(r'\s+', ' ', text)
    words = [w.strip() for w in text.split() if w.strip()]
    return words

def process_slide_text(slide_text_lines):
    """
    Detecta automáticamente:
    - //Te adoro a Ti// → REPITE_ULTIMA_FRASE:2
    - //Todo el slide// → DUPLICADO
    - Frase al final con // → REPITE_ULTIMA_FRASE
    """
    full_text = " ".join(slide_text_lines).strip()
    print(f"Procesando slide: '{full_text}'")
    
    metadata = []
    processed_words = []
    raw_lines = slide_text_lines.copy()

    # Caso 1: Todo el slide entre // → duplicar todo el slide
    if full_text.startswith("//") and full_text.endswith("//"):
        content = full_text[2:-2].strip()
        print("→ Todo el slide entre // → DUPLICADO")
        words = clean_and_tokenize(content)
        processed_words = words + words
        metadata.append("DUPLICADO")
        metadata.append(f"🔄MITAD1:{len(words)}")
        return {"text": processed_words, "metadata": metadata, "raw_text": raw_lines}

    # Caso 2: Hay // dentro del texto → analizar dónde está
    if "//" in full_text:
        parts = [p.strip() for p in full_text.split("//") if p.strip()]
        
        if len(parts) >= 2:
            before = parts[0]
            to_repeat = parts[1]
            after = " ".join(parts[2:]) if len(parts) > 2 else ""

            # Subcaso A: La frase a repetir está al FINAL → REPITE_ULTIMA_FRASE
            if full_text.strip().endswith("//" + to_repeat + "//") or full_text.strip().endswith(to_repeat + "//"):
                print(f"→ REPITE_ULTIMA_FRASE detectada: '{to_repeat}'")
                words_before = clean_and_tokenize(before + " " + after)
                words_repeat = clean_and_tokenize(to_repeat)
                processed_words = words_before + words_repeat + words_repeat  # repite 2 veces por defecto
                metadata.append(f"REPITE_ULTIMA_FRASE:2")
                metadata.append(f"FRASE_REPETIDA: {' '.join(words_repeat)}")
                
            # Subcaso B: La frase a repetir está en medio o al inicio → DUPLICADO clásico
            else:
                print("→ // en medio → DUPLICADO clásico")
                all_content = before + " " + to_repeat
                words = clean_and_tokenize(all_content)
                processed_words = words + words
                metadata.append("DUPLICADO")
                metadata.append(f"🔄MITAD1:{len(words)}")

            return {"text": processed_words, "metadata": metadata, "raw_text": raw_lines}

    # Caso 3: Sin // → slide normal
    processed_words = clean_and_tokenize(full_text)
    return {"text": processed_words, "metadata": metadata, "raw_text": raw_lines}

def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        slides_data = {}
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Separar por líneas para detectar mejor
                    lines = [line.strip() for line in shape.text.split('\n') if line.strip()]
                    slide_text_lines.extend(lines)
            
            if slide_text_lines:
                result = process_slide_text(slide_text_lines)
                slides_data[f"slide_{i}"] = {
                    "raw_text": result["raw_text"],
                    "processed_text": result["text"],
                    "metadata": result["metadata"]
                }
                print(f"  Slide {i} → {len(result['text'])} palabras | Metadata: {result['metadata']}")
        
        return slides_data
    except Exception as e:
        print(f"Error al procesar {pptx_path}: {e}")
        return {}

def save_to_json(data, output_path):
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2, separators=(',', ': '))

def main():
    input_dir = "test_files"
    output_dir = "."
    
    if not os.path.exists(input_dir):
        print(f"No se encuentra la carpeta: {input_dir}")
        return
        
    pptx_files = glob.glob(os.path.join(input_dir, "*.pptx"))
    
    if not pptx_files:
        print(f"No hay archivos .pptx en {input_dir}")
        return
    
    print(f"Encontrados {len(pptx_files)} archivos para procesar:\n")
    
    for pptx_file in pptx_files:
        filename = os.path.basename(pptx_file)
        print(f"Procesando → {filename}")
        
        slides_data = extract_text_from_pptx(pptx_file)
        
        if slides_data:
            output_file = os.path.join(output_dir, os.path.splitext(filename)[0] + "_lyrics.json")
            save_to_json(slides_data, output_file)
            print(f"Guardado: {output_file} ({len(slides_data)} slides)\n")
        else:
            print(f"No se pudo procesar {filename}\n")

if __name__ == "__main__":
    main()